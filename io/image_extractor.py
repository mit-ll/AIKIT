from pptx import Presentation
import os
import copy
import fitz 
import io
from typing import AsyncIterator, Iterator
from PIL import Image, UnidentifiedImageError
#from wand.image import Image as WandImage


class ImageExtractor:
    def __init__(self):
        pass
    
    def extract_images_from_pdf(self, pdf_path, output_folder):
        extracted_files = []
        metadata_list = []

        # Open the PDF file
        pdf_document = fitz.open(pdf_path)

        base_name = os.path.basename(pdf_path)
        name, _ = os.path.splitext(base_name)
        
        # Iterate through each page
        for page_number in range(len(pdf_document)):
            page = pdf_document.load_page(page_number)
            images = page.get_images(full=True)
            
            for img_index, img in enumerate(images):
                if img_index > 100:
                    print('More > 100 images on a single page. Skipping remainder and going to next page.')
                    break
                xref = img[0]
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                
                pix = fitz.Pixmap(pdf_document, xref)
                image_path = os.path.join(output_folder, f"{name}_image_page{page_number+1}_{img_index}.{image_ext}")
                if pix.width > 100 and  pix.height > 100:

                    try:
                        image = Image.open(io.BytesIO(image_bytes))
                        extracted_files.append( image_path )
                        metadata_list.append({'pdf_path':pdf_path, 'source_file_name':name, 'page_number':page_number+1, 'img_index':img_index})

                        # Convert the image to RGB format to ensure it's not grayscale
                        if image.mode != 'RGB':
                            image = image.convert('RGB')

                        # Save the image as a file
                        with open(image_path, "wb") as image_file:
                            image.save(image_file, format='PNG')
                    except UnidentifiedImageError:
                        print(f"Cannot identify image file. Please check the source. {image_path} pg {page_number+1}")

                else:
                    print(f"Skipping image {image_path} with width < 1 pixel on page {page_number+1}")

                    
        print("PDF Image extraction completed.")
        print(extracted_files)
        return extracted_files, metadata_list

    def extract_images_and_captions_from_pptx(self, pptx_path, output_dir):

        extracted_files, metadata_list = [], []
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        presentation = Presentation(pptx_path)
        base_name = os.path.basename(pptx_path)
        # Split the base name and extension
        ppt_name, _ = os.path.splitext(base_name)
        image_count = 0
        extracted_files = [] 
        for slide_ix, slide in enumerate(presentation.slides):
            for shape_ix, shape in enumerate(slide.shapes):
                if shape.shape_type == 13:  # Shape type 13 is a picture
                    image = shape.image
                    image_bytes = image.blob
                    image_format = image.ext

                    output_path = os.path.join(output_dir, f'{ppt_name}_image_page_{slide_ix+1}_{image_count}.png')
                    # Ensure the image format is PNG, if not, convert it
                    
                    # ignore windows shapes?
                    if image_format == 'wmf':
                        # Convert WMF to PNG using Wand, (Windows Metafile, shapes)
                        #with WandImage(blob=image_bytes, format="wmf") as img:
                        #    img.format = 'png'
                        #    img.save(filename=output_path)
                        pass
                    elif image_format != 'png':
                        print(f'png conversion {ppt_name} format {image_format}')
                        with Image.open(io.BytesIO(image_bytes)) as img:
                            img.save(output_path, format='PNG')
                    else:
                        with open(output_path, 'wb') as f:
                            f.write(image_bytes)

                    # Check for caption (text in a textbox near the image)
                    caption = ""
                    for shape_ in slide.shapes:
                        if shape_ != shape and hasattr(shape_, "text"):
                            if shape_.left > shape.left and shape_.left < (shape.left + shape.width):
                                if shape_.top > (shape.top + shape.height) and shape_.top < (shape.top + shape.height + 200000):  # Adjust this range as needed
                                    caption = shape_.text


                    image_count += 1
                    #ignore tiny images and wmf files
                    if shape.width >= 100 and shape.height >=100 and image_format != 'wmf':
                        extracted_files.append(output_path)
                        metadata_list.append( {'pptx_path':pptx_path, 'ppt_name':ppt_name, 'slide_ix':slide_ix, 'shape_ix':shape_ix, 'caption':caption} )

        print(f"Extracted {image_count} images with captions to {output_dir}")
        return extracted_files, metadata_list

    def extract(self, file_path, output_folder):
        if file_path.lower().endswith('.pptx') or file_path.lower().endswith('.ppt'):
            return self.extract_images_and_captions_from_pptx(file_path, output_folder)
        elif file_path.lower().endswith('.pdf'):
            return self.extract_images_from_pdf(file_path, output_folder)
        else:
            raise ValueError("Unsupported file type or path provided")

# Example usage
if __name__ == "__main__":
    pptx_path = "/home/da31742/AIKIT/io/afmed_docs/r2ai_wip.pptx"
    output_dir = "/home/da31742/AIKIT/io/output_images"
    ImageExtractor().extract(pptx_path, output_dir)#extract_images_and_captions_from_pptx(pptx_path, output_dir)

